# ============================================================
# MedStudyFlow — Aplicación de flashcards médicas interactivas
# Autor: Proyecto educativo
# Versión: 1.0 (compatible con Streamlit Cloud)
# ============================================================

import streamlit as st
import pandas as pd
import sqlite3
import fitz  # PyMuPDF para PDFs
from pptx import Presentation
from PIL import Image
import io
import openai
import os

# ------------------------- CONFIGURACIÓN GENERAL -------------------------
st.set_page_config(
    page_title="MedStudyFlow",
    page_icon="🧠",
    layout="centered"
)

st.markdown("""
    <style>
    body {background-color:#fefefe;}
    .main {background-color:#faf9f9;}
    h1,h2,h3 {color:#4A4A4A;font-family:'Nunito Sans',sans-serif;}
    .stButton>button {
        background-color:#F5A6C1;
        color:white;
        border:none;
        border-radius:12px;
        padding:8px 16px;
        font-weight:600;
    }
    .stProgress > div > div > div > div {
        background-color:#F5A6C1;
    }
    </style>
""", unsafe_allow_html=True)

st.title("💗 MedStudyFlow")
st.write("Tu asistente visual de estudio médico — crea, valida y practica con tus propias flashcards.")

# ------------------------- API KEY -------------------------
with st.expander("🔐 Configuración de clave API"):
    st.markdown("Introduce tu clave de OpenAI (no se guarda):")
    api_key = st.text_input("Clave OpenAI:", type="password")
    if api_key:
        openai.api_key = api_key

# ------------------------- BASE DE DATOS LOCAL -------------------------
os.makedirs("data", exist_ok=True)
conn = sqlite3.connect("data/resultados_medstudyflow.db")
cursor = conn.cursor()
cursor.execute("""
CREATE TABLE IF NOT EXISTS resultados (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    tema TEXT,
    aciertos INTEGER,
    total INTEGER,
    nivel TEXT
)
""")
conn.commit()

# ------------------------- FUNCIÓN: EXTRAER TEXTO -------------------------
def extraer_texto(archivo):
    nombre = archivo.name.lower()
    texto = ""
    if nombre.endswith(".pdf"):
        doc = fitz.open(stream=archivo.read(), filetype="pdf")
        for page in doc:
            texto += page.get_text()
    elif nombre.endswith(".pptx"):
        prs = Presentation(io.BytesIO(archivo.read()))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
    elif nombre.endswith((".txt", ".csv", ".xlsx")):
        if nombre.endswith(".csv"):
            df = pd.read_csv(archivo)
        elif nombre.endswith(".xlsx"):
            df = pd.read_excel(archivo)
        else:
            df = pd.DataFrame({"texto": archivo.read().decode("utf-8").splitlines()})
        texto += " ".join(df.astype(str).values.flatten())
    elif nombre.endswith((".jpg", ".png", ".jpeg")):
        texto += "[Imagen cargada: el análisis visual no está habilitado aquí.]"
    else:
        texto = archivo.read().decode(errors="ignore")
    return texto.strip()

# ------------------------- FUNCIÓN: VERIFICACIÓN MÉDICA -------------------------
def verificar_contenido(texto):
    prompt = f"""
Eres un profesor de medicina. Evalúa la precisión científica del siguiente texto:
---
{texto[:2000]}
---
Devuelve un resumen con:
1. Nivel de corrección (Correcto, Parcial o Incorrecto)
2. Breve explicación (máx 2 líneas)
3. Referencia o corrección sugerida.
"""
    try:
        respuesta = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[{"role": "system", "content": "Eres un médico docente preciso y claro."},
                      {"role": "user", "content": prompt}],
            max_tokens=300
        )
        return respuesta.choices[0].message.content.strip()
    except Exception as e:
        return f"⚠️ Error al verificar: {e}"

# ------------------------- FUNCIÓN: GENERAR PREGUNTAS -------------------------
def generar_preguntas(texto, tipo="basico", n=5):
    prompt = f"""
Crea {n} preguntas tipo MIR/USMLE sobre el siguiente contenido:
---
{texto[:2500]}
---
Nivel: {'básico (anatomía, bioquímica, fisiología)' if tipo=='basico' else 'clínico (semiología, pediatría, cirugía, neurología)'}
Cada pregunta debe incluir:
- Enunciado claro
- 4 opciones (A–D)
- Una respuesta correcta
- Explicación breve y colorida
Formato:
PREGUNTA: ...
A) ...
B) ...
C) ...
D) ...
RESPUESTA: ...
EXPLICACIÓN: ...
"""
    try:
        resp = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[{"role": "system", "content": "Eres un profesor de medicina que formula preguntas didácticas."},
                      {"role": "user", "content": prompt}],
            max_tokens=1200
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"⚠️ Error al generar preguntas: {e}"

# ------------------------- INTERFAZ DE USUARIO -------------------------
st.header("📂 1. Sube tus archivos de estudio")
archivos = st.file_uploader(
    "Sube tus flashcards, presentaciones o notas",
    type=["pdf", "pptx", "txt", "csv", "xlsx", "jpg", "png"],
    accept_multiple_files=True
)

if archivos and api_key:
    contenido_total = ""
    for archivo in archivos:
        st.write(f"📘 Procesando: `{archivo.name}` ...")
        texto = extraer_texto(archivo)
        contenido_total += texto + "\n"
        st.success(f"✅ Texto extraído: {len(texto.split())} palabras")

    # Verificación médica
    st.header("🩺 2. Verificación médica del contenido")
    if st.button("Verificar precisión médica"):
        with st.spinner("Analizando contenido..."):
            verificacion = verificar_contenido(contenido_total)
        st.info(verificacion)

    # Generación de preguntas
    st.header("🧠 3. Generación de preguntas tipo examen")
    tipo_examen = st.radio("Selecciona tipo de examen:", ["Básico", "Clínico"])
    n_preguntas = st.slider("Número de preguntas", 3, 10, 5)

    if st.button("Generar examen"):
        with st.spinner("Creando preguntas..."):
            preguntas = generar_preguntas(
                contenido_total,
                tipo="basico" if tipo_examen == "Básico" else "clinico",
                n=n_preguntas
            )
        st.markdown("### 🧾 Preguntas generadas")
        st.text_area("Preguntas", preguntas, height=400)

        # Guarda el resultado localmente
        cursor.execute(
            "INSERT INTO resultados (tema, aciertos, total, nivel) VALUES (?,?,?,?)",
            ("Autoexamen", 0, n_preguntas, tipo_examen)
        )
        conn.commit()
        st.success("✅ Examen registrado en tu progreso local.")

# ------------------------- HISTORIAL DE PROGRESO -------------------------
st.header("📊 4. Historial de progreso")
df = pd.read_sql_query("SELECT * FROM resultados", conn)
if not df.empty:
    df["porcentaje"] = (df["aciertos"] / df["total"] * 100).fillna(0)
    st.dataframe(df[["tema", "aciertos", "total", "nivel", "porcentaje"]])
else:
    st.info("Aún no hay resultados guardados.")

st.write("---")
st.caption("© 2025 MedStudyFlow — Aplicación educativa desarrollada con ❤️ y GPT-4-turbo.")
